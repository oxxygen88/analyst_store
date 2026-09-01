from __future__ import annotations

import json
import re
import urllib.error
import urllib.parse
import urllib.request
import importlib.util
from io import BytesIO
from typing import Any, Dict

import numpy as np
import pandas as pd

from .analytics import (
    anomaly_tables,
    commercial_kpis,
    current_stock_snapshot,
    filter_period,
    inventory_health,
    monthly_sales,
    opportunity_scoring,
    pareto_products,
    revenue_inventory_matrix,
    target_status,
)
from .utils import pct, rupiah
from .gemini_models import call_gemini_with_fallback


def _jsonable(value: Any):
    if value is None:
        return None
    if isinstance(value, (np.integer,)):
        return int(value)
    if isinstance(value, (np.floating,)):
        if np.isnan(value) or np.isinf(value):
            return None
        return float(value)
    if isinstance(value, pd.Timestamp):
        return value.isoformat()
    if isinstance(value, float) and (np.isnan(value) or np.isinf(value)):
        return None
    return value


def _records(df: pd.DataFrame, columns: list[str], n: int = 10) -> list[dict]:
    if df is None or df.empty:
        return []
    cols = [c for c in columns if c in df.columns]
    out = df[cols].head(n).copy()
    records = []
    for row in out.to_dict("records"):
        records.append({k: _jsonable(v) for k, v in row.items()})
    return records


def presentation_context(bundle, as_of: pd.Timestamp, location: str = "IDK-ATP") -> Dict[str, Any]:
    """Compact, aggregated context safe to send to AI. Raw transaction rows are intentionally excluded."""
    as_of = pd.Timestamp(as_of).normalize()
    start = as_of.to_period("M").to_timestamp()
    tx_period = filter_period(bundle.tx, start, as_of)
    kpi = commercial_kpis(tx_period)
    inv = inventory_health(bundle.opening, bundle.tx, as_of, bundle.purchases)
    p = opportunity_scoring(pareto_products(bundle.tx, inv, start, as_of), bundle.tx, as_of)
    target = target_status(bundle.tx, bundle.targets, start, as_of, location)
    monthly = monthly_sales(bundle.tx)
    monthly = monthly[monthly["month"].le(as_of.to_period("M").to_timestamp())].copy()
    if bundle.targets is not None and not bundle.targets.empty:
        monthly = monthly.merge(
            bundle.targets[["bulan", "target_omzet"]].rename(columns={"bulan": "month"}),
            on="month",
            how="left",
        )
    else:
        monthly["target_omzet"] = np.nan

    inv_summary = (
        inv.groupby("inventory_status", as_index=False)
        .agg(sku_count=("sku", "nunique"), stock_qty=("current_stock", "sum"), stock_value=("current_stock_value", lambda s: float(s.clip(lower=0).sum())))
        .sort_values("stock_value", ascending=False)
    )

    core = p[p["pareto_group"].eq("CORE_20")] if not p.empty else pd.DataFrame()
    opp = p[p["pareto_group"].eq("OPPORTUNITY")] if not p.empty else pd.DataFrame()
    longtail = p[p["pareto_group"].eq("LONG_TAIL")] if not p.empty else pd.DataFrame()
    total_rev = float(p["revenue"].sum()) if not p.empty else 0.0
    core_share = float(core["revenue"].sum() / total_rev) if total_rev else 0.0
    opp_share = float(opp["revenue"].sum() / total_rev) if total_rev else 0.0
    a80_count = int(p["a80_member"].sum()) if not p.empty else 0

    focus = p[p["recommended_action"].isin([
        "Push sales", "Push / campaign candidate", "Replenish before push",
        "Protect sales / replenish", "Emergency replenish",
    ])].copy() if not p.empty else pd.DataFrame()
    if not focus.empty:
        focus["priority"] = focus["recommended_action"].map({
            "Emergency replenish": 1,
            "Protect sales / replenish": 2,
            "Push sales": 3,
            "Push / campaign candidate": 4,
            "Replenish before push": 5,
        }).fillna(9)
        focus = focus.sort_values(["priority", "opportunity_score", "revenue"], ascending=[True, False, False])

    supplier = revenue_inventory_matrix(bundle.tx, inv, start, as_of, "supplier")
    category = revenue_inventory_matrix(bundle.tx, inv, start, as_of, "subdept")
    issues = anomaly_tables(bundle.opening, bundle.tx[bundle.tx["date"].le(as_of)], inv)

    target_dict = None
    if target is not None:
        target_dict = {
            "target": target.target,
            "actual": target.actual,
            "achievement": target.achievement,
            "gap": target.gap,
            "daily_run_rate": target.daily_run_rate,
            "required_daily_sales": target.required_daily_sales,
            "projected_month_end": target.projected_month_end,
            "projected_gap": target.projected_gap,
            "pace_achievement": target.pace_achievement,
            "remaining_days": target.remaining_days,
        }

    return {
        "branch": location,
        "as_of": as_of.strftime("%Y-%m-%d"),
        "period": {"start": start.strftime("%Y-%m-%d"), "end": as_of.strftime("%Y-%m-%d")},
        "kpi": {k: _jsonable(v) for k, v in kpi.items()},
        "target": target_dict,
        "pareto": {
            "active_selling_sku": int(len(p)),
            "core20_sku": int(len(core)),
            "core_share": core_share,
            "opportunity_sku": int(len(opp)),
            "opportunity_share": opp_share,
            "long_tail_sku": int(len(longtail)),
            "a80_sku": a80_count,
            "core_stockout": int((core["current_stock"].le(0)).sum()) if not core.empty else 0,
        },
        "monthly": _records(monthly, ["month", "net_sales", "target_omzet", "trx_count", "atv", "gross_profit", "gross_margin"], 24),
        "top_focus_products": _records(focus, ["sku", "nama_barang", "supplier", "pareto_group", "revenue", "revenue_share", "growth_30d", "gross_margin", "current_stock", "stock_cover_days", "inventory_status", "opportunity_score", "recommended_action"], 15),
        "inventory_summary": _records(inv_summary, ["inventory_status", "sku_count", "stock_qty", "stock_value"], 20),
        "top_suppliers": _records(supplier, ["supplier", "revenue", "inventory_value", "revenue_share", "inventory_share", "productivity_index"], 10),
        "top_categories": _records(category, ["subdept", "revenue", "inventory_value", "revenue_share", "inventory_share", "productivity_index"], 10),
        "anomalies": {name: int(len(df)) for name, df in issues.items()},
    }


def _extract_json(text: str) -> dict:
    raw = (text or "").strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.I)
        raw = re.sub(r"\s*```$", "", raw)
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r"\{.*\}", raw, flags=re.S)
        if not m:
            raise ValueError("Respons AI tidak berisi JSON yang dapat dibaca.")
        return json.loads(m.group(0))


def _module_available(name: str) -> bool:
    try:
        return importlib.util.find_spec(name) is not None
    except (ModuleNotFoundError, ValueError):
        return False


def ai_runtime_info() -> Dict[str, bool]:
    """Return dependency status without importing heavy packages."""
    return {
        "openai_sdk": _module_available("openai"),
        "gemini_sdk": _module_available("google.genai"),
        "python_pptx": _module_available("pptx"),
    }


def _responses_output_text(payload: Dict[str, Any]) -> str:
    """Extract text from the raw JSON returned by the Responses API."""
    chunks = []
    for item in payload.get("output", []) or []:
        if not isinstance(item, dict):
            continue
        for content in item.get("content", []) or []:
            if not isinstance(content, dict):
                continue
            if content.get("type") in {"output_text", "text"} and content.get("text"):
                chunks.append(str(content["text"]))
    text = "\n".join(chunks).strip()
    if not text and payload.get("output_text"):
        text = str(payload["output_text"]).strip()
    if not text:
        raise RuntimeError("OpenAI mengembalikan respons tanpa teks yang dapat dibaca.")
    return text


def _call_openai_responses_http(api_key: str, model: str, prompt: str) -> str:
    """Dependency-free HTTPS fallback when the openai Python SDK is unavailable."""
    body = json.dumps({"model": model, "input": prompt}, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        "https://api.openai.com/v1/responses",
        data=body,
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            payload = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        raw = exc.read().decode("utf-8", errors="replace")
        try:
            message = json.loads(raw).get("error", {}).get("message") or raw
        except Exception:
            message = raw
        if exc.code == 401:
            raise RuntimeError("OpenAI API key ditolak (401). Periksa API key dan project API Anda.") from exc
        if exc.code == 429:
            raise RuntimeError(f"OpenAI API menolak request karena limit/billing (429): {message}") from exc
        raise RuntimeError(f"OpenAI API error HTTP {exc.code}: {message}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"Tidak dapat terhubung ke OpenAI API: {exc.reason}") from exc
    return _responses_output_text(payload)



def _call_gemini_http(api_key: str, model: str, prompt: str) -> str:
    """Gemini REST call with stable-model auto fallback."""
    return call_gemini_with_fallback(
        api_key, model, prompt, json_mode=True, prefer_sdk=False
    )


def _call_gemini_sdk(api_key: str, model: str, prompt: str) -> str:
    """Gemini SDK call with stable-model auto fallback."""
    return call_gemini_with_fallback(
        api_key, model, prompt, json_mode=True, prefer_sdk=True
    )

def generate_ai_insights(api_key: str, model: str, context: Dict[str, Any], *, audience: str, language: str, provider: str = "openai") -> Dict[str, Any]:
    """Generate structured management insight using OpenAI or Google Gemini."""
    provider_key = (provider or "openai").strip().lower()
    if provider_key not in {"openai", "gemini"}:
        raise ValueError(f"AI provider tidak dikenal: {provider}")
    if not api_key:
        provider_name = "Gemini" if provider_key == "gemini" else "OpenAI"
        raise ValueError(f"{provider_name} API key belum diisi.")

    lang_instruction = "Bahasa Indonesia yang profesional dan ringkas" if language == "Bahasa Indonesia" else "professional concise English"
    schema = {
        "presentation_title": "string",
        "executive_summary": ["3-5 concise strings"],
        "target_insights": ["2-4 strings"],
        "pareto_insights": ["2-4 strings"],
        "inventory_insights": ["2-4 strings"],
        "profitability_insights": ["1-3 strings"],
        "risks": ["2-4 strings"],
        "recommended_actions": [
            {"priority": 1, "action": "string", "why": "string", "expected_impact": "string", "owner": "string"}
        ],
        "closing_message": "string",
    }
    prompt = f"""
You are a senior retail performance analyst preparing a management presentation for a Mom & Baby retail branch.
Audience: {audience}
Write in: {lang_instruction}.

STRICT DATA RULES:
- Use ONLY the numbers and facts in CONTEXT below.
- Do not invent revenue, margin, targets, product performance, or causal explanations.
- If causality is not proven, phrase it as a hypothesis/recommendation, not a fact.
- Prioritize actions that help close the monthly target while protecting Core products and avoiding unhealthy inventory.
- Treat the Pareto result dynamically: Core 20% is not automatically equal to 80% revenue.
- HPP/Gross Profit are estimates when present.
- Focus on management-ready insights, not generic advice.
- Do not mention that you are an AI.

Return ONLY valid JSON matching this shape:
{json.dumps(schema, ensure_ascii=False)}

CONTEXT:
{json.dumps(context, ensure_ascii=False, default=_jsonable)}
""".strip()

    if provider_key == "gemini":
        if _module_available("google.genai"):
            response_text = _call_gemini_sdk(api_key, model, prompt)
        else:
            response_text = _call_gemini_http(api_key, model, prompt)
    else:
        # Prefer the official SDK, but do not make the entire AI menu depend on it.
        try:
            from openai import OpenAI
        except ImportError:
            response_text = _call_openai_responses_http(api_key, model, prompt)
        else:
            try:
                client = OpenAI(api_key=api_key, timeout=120.0)
                response = client.responses.create(model=model, input=prompt)
                response_text = response.output_text
            except Exception as exc:
                raise RuntimeError(f"OpenAI API request gagal: {exc}") from exc

    data = _extract_json(response_text)

    # Defensive normalization so the PPT builder always receives predictable structures.
    for key in ["executive_summary", "target_insights", "pareto_insights", "inventory_insights", "profitability_insights", "risks"]:
        val = data.get(key, [])
        if not isinstance(val, list):
            val = [str(val)]
        data[key] = [str(x) for x in val if str(x).strip()][:6]
    actions = data.get("recommended_actions", [])
    if not isinstance(actions, list):
        actions = []
    norm_actions = []
    for i, row in enumerate(actions[:8], start=1):
        if isinstance(row, dict):
            norm_actions.append({
                "priority": row.get("priority", i),
                "action": str(row.get("action", "")),
                "why": str(row.get("why", "")),
                "expected_impact": str(row.get("expected_impact", "")),
                "owner": str(row.get("owner", "Management")),
            })
    data["recommended_actions"] = norm_actions
    return data


# ---------- PowerPoint builder ----------

def build_pptx(context: Dict[str, Any], ai: Dict[str, Any], *, branch: str, as_of: pd.Timestamp, audience: str) -> bytes:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("Package 'python-pptx' belum terinstall. Jalankan setup_windows.bat kembali setelah update aplikasi.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BLUE = RGBColor(47, 104, 176)
    ORANGE = RGBColor(246, 139, 36)
    DARK = RGBColor(31, 41, 55)
    MUTED = RGBColor(100, 116, 139)
    LIGHT = RGBColor(244, 247, 251)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(22, 163, 74)
    RED = RGBColor(220, 38, 38)

    def add_rect(slide, x, y, w, h, fill, radius=False):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None, valign=None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.clear(); tf.word_wrap = True
        if valign is not None: tf.vertical_anchor = valign
        p = tf.paragraphs[0]; p.text = str(text)
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
        if align is not None: p.alignment = align
        return box

    def slide_header(slide, title, subtitle=None, num=None):
        add_text(slide, title, .7, .35, 11.7, .5, 25, DARK, True)
        if subtitle:
            add_text(slide, subtitle, .72, .9, 11.7, .32, 10.5, MUTED)
        add_rect(slide, .7, 1.26, 1.0, .055, ORANGE)
        if num is not None:
            add_text(slide, str(num), 12.15, .38, .45, .3, 9, MUTED, False, PP_ALIGN.RIGHT)

    def footer(slide):
        add_text(slide, f"INDOKIDS · {branch} · Data s.d. {pd.Timestamp(as_of):%d %b %Y}", .7, 7.12, 11.8, .22, 8.5, MUTED)

    def bullet_block(slide, items, x, y, w, h, size=16):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.clear(); tf.word_wrap = True
        for i, item in enumerate(items or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item); p.level = 0
            p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = DARK
            p.space_after = Pt(9)
            p.text = "• " + p.text
        return box

    def metric_card(slide, x, y, w, title, value, note=""):
        add_rect(slide, x, y, w, 1.08, LIGHT, True)
        add_text(slide, title, x+.16, y+.13, w-.32, .22, 9.5, MUTED, True)
        add_text(slide, value, x+.16, y+.42, w-.32, .32, 19, DARK, True)
        if note:
            add_text(slide, note, x+.16, y+.79, w-.32, .18, 8.5, MUTED)

    # 1 Title
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.333, 7.5, BLUE)
    add_rect(slide, .72, 1.45, .12, 2.25, ORANGE)
    add_text(slide, ai.get("presentation_title") or "Branch Performance & Target Recovery", 1.1, 1.45, 10.8, 1.3, 32, WHITE, True)
    add_text(slide, f"{branch} · Management Brief", 1.1, 2.9, 9.6, .5, 19, WHITE)
    add_text(slide, f"Data sampai {pd.Timestamp(as_of):%d %B %Y} · Audience: {audience}", 1.1, 3.52, 9.8, .35, 11, WHITE)
    add_text(slide, "Monitor · Diagnose · Act", 1.1, 5.85, 4.8, .4, 13, WHITE, True)

    # 2 Executive overview
    slide = prs.slides.add_slide(blank); slide_header(slide, "Executive Overview", "Current performance, target pace, and commercial health", 2)
    k = context["kpi"]; t = context.get("target") or {}
    metric_card(slide, .7, 1.55, 2.85, "NET SALES MTD", rupiah(k.get("net_sales")))
    metric_card(slide, 3.75, 1.55, 2.85, "TARGET", rupiah(t.get("target")) if t else "-")
    metric_card(slide, 6.8, 1.55, 2.85, "ACHIEVEMENT", pct(t.get("achievement")) if t else "-")
    metric_card(slide, 9.85, 1.55, 2.78, "PROJECTED GAP", rupiah(t.get("projected_gap")) if t else "-")
    metric_card(slide, .7, 2.85, 2.85, "TRANSACTIONS", f"{int(k.get('trx_count',0)):,}".replace(",", "."))
    metric_card(slide, 3.75, 2.85, 2.85, "AVG TRANSACTION", rupiah(k.get("atv")))
    metric_card(slide, 6.8, 2.85, 2.85, "EST. GROSS PROFIT", rupiah(k.get("gross_profit")))
    metric_card(slide, 9.85, 2.85, 2.78, "GROSS MARGIN", pct(k.get("gross_margin")))
    add_text(slide, "AI readout", .72, 4.28, 2.0, .32, 13, BLUE, True)
    bullet_block(slide, ai.get("executive_summary", []), .8, 4.72, 11.7, 1.95, 15)
    footer(slide)

    # 3 Target trend
    slide = prs.slides.add_slide(blank); slide_header(slide, "Target Performance", "Actual sales vs monthly target", 3)
    monthly = context.get("monthly", [])
    if monthly:
        chart_data = CategoryChartData()
        chart_data.categories = [pd.Timestamp(r["month"]).strftime("%b") for r in monthly]
        chart_data.add_series("Net Sales", [float(r.get("net_sales") or 0) for r in monthly])
        chart_data.add_series("Target", [float(r.get("target_omzet") or 0) for r in monthly])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.75), Inches(1.6), Inches(7.55), Inches(4.85), chart_data).chart
        chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.value_axis.tick_labels.number_format = '0,,"M"'
        chart.value_axis.has_major_gridlines = True
        chart.chart_title.text_frame.text = ""
    add_text(slide, "What matters now", 8.7, 1.7, 3.6, .35, 14, BLUE, True)
    bullet_block(slide, ai.get("target_insights", []), 8.7, 2.2, 3.8, 3.8, 14)
    if t:
        add_text(slide, f"Required pace: {rupiah(t.get('required_daily_sales'))}/hari", 8.7, 6.1, 3.8, .3, 11, RED if (t.get("projected_gap") or 0) > 0 else GREEN, True)
    footer(slide)

    # 4 Pareto
    slide = prs.slides.add_slide(blank); slide_header(slide, "Pareto & Revenue Concentration", "Protect winners, then grow the opportunity layer", 4)
    par = context["pareto"]
    chart_data = CategoryChartData(); chart_data.categories = ["Core 20", "Opportunity", "Long Tail"]
    chart_data.add_series("Revenue Share", [par.get("core_share",0)*100, par.get("opportunity_share",0)*100, max(0,1-par.get("core_share",0)-par.get("opportunity_share",0))*100])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(.8), Inches(1.7), Inches(5.4), Inches(3.9), chart_data).chart
    chart.has_legend = False; chart.value_axis.tick_labels.number_format = '0"%"'; chart.value_axis.maximum_scale = 100
    metric_card(slide, .85, 5.8, 2.4, "CORE 20 SHARE", pct(par.get("core_share")))
    metric_card(slide, 3.45, 5.8, 2.4, "A80 SKU", f"{par.get('a80_sku',0):,}".replace(",", "."))
    add_text(slide, "AI insights", 6.65, 1.75, 2.2, .35, 14, BLUE, True)
    bullet_block(slide, ai.get("pareto_insights", []), 6.65, 2.2, 5.8, 3.8, 14)
    footer(slide)

    # 5 Product focus table
    slide = prs.slides.add_slide(blank); slide_header(slide, "Products to Improve", "Highest-priority growth and availability actions", 5)
    rows = context.get("top_focus_products", [])[:8]
    headers = ["SKU", "Product", "Revenue", "Growth", "Stock", "Status", "Action"]
    table = slide.shapes.add_table(len(rows)+1, len(headers), Inches(.6), Inches(1.55), Inches(12.15), Inches(4.9)).table
    widths = [1.25, 3.1, 1.45, 1.0, .75, 1.25, 3.35]
    for i,w in enumerate(widths): table.columns[i].width = Inches(w)
    for c,h in enumerate(headers):
        cell=table.cell(0,c); cell.text=h; cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
        for p in cell.text_frame.paragraphs: p.font.color.rgb=WHITE; p.font.bold=True; p.font.size=Pt(10)
    for r_idx,row in enumerate(rows,1):
        vals=[row.get("sku",""), str(row.get("nama_barang","") or "")[:34], rupiah(row.get("revenue")), pct(row.get("growth_30d")), f"{float(row.get('current_stock') or 0):,.0f}".replace(",","."), row.get("inventory_status",""), row.get("recommended_action","")]
        for c,val in enumerate(vals):
            cell=table.cell(r_idx,c); cell.text=str(val); cell.fill.solid(); cell.fill.fore_color.rgb=WHITE if r_idx%2 else LIGHT
            for p in cell.text_frame.paragraphs: p.font.size=Pt(8.5); p.font.color.rgb=DARK
    footer(slide)

    # 6 Inventory
    slide = prs.slides.add_slide(blank); slide_header(slide, "Inventory Health", "Revenue protection requires healthy availability and controlled capital", 6)
    invsum = context.get("inventory_summary", [])
    if invsum:
        chart_data = CategoryChartData(); chart_data.categories=[str(r.get("inventory_status","")) for r in invsum]
        chart_data.add_series("SKU", [int(r.get("sku_count") or 0) for r in invsum])
        chart=slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(.75), Inches(1.65), Inches(7.0), Inches(4.65), chart_data).chart
        chart.has_legend=False; chart.value_axis.has_major_gridlines=True
    add_text(slide, "AI inventory readout", 8.15, 1.72, 3.4, .35, 14, BLUE, True)
    bullet_block(slide, ai.get("inventory_insights", []), 8.15, 2.18, 4.3, 2.55, 14)
    add_text(slide, "Risks", 8.15, 4.9, 2.0, .3, 13, RED, True)
    bullet_block(slide, ai.get("risks", []), 8.15, 5.25, 4.3, 1.25, 12.5)
    footer(slide)

    # 7 Action plan
    slide = prs.slides.add_slide(blank); slide_header(slide, "Recommended Actions", "Prioritized actions generated from the current branch data", 7)
    actions = ai.get("recommended_actions", [])[:6]
    y=1.55
    for i,a in enumerate(actions,1):
        add_rect(slide,.75,y,11.85,.72,LIGHT,True)
        add_text(slide,str(a.get("priority",i)),.95,y+.16,.35,.26,12,BLUE,True,PP_ALIGN.CENTER)
        add_text(slide,a.get("action",""),1.4,y+.10,4.25,.25,12,DARK,True)
        add_text(slide,a.get("why",""),1.4,y+.38,5.8,.2,8.5,MUTED)
        add_text(slide,a.get("expected_impact",""),7.35,y+.14,3.2,.2,9.5,DARK,True)
        add_text(slide,a.get("owner",""),10.75,y+.14,1.4,.2,9,MUTED,True)
        y += .83
    footer(slide)

    # 8 Closing
    slide = prs.slides.add_slide(blank)
    add_rect(slide,0,0,13.333,7.5,DARK)
    add_text(slide,"Management Focus",.85,1.05,5.5,.6,29,WHITE,True)
    add_rect(slide,.87,1.85,1.0,.06,ORANGE)
    add_text(slide,ai.get("closing_message") or "Protect winners. Grow potentials. Fix availability. Reduce waste.",.87,2.4,10.8,1.4,25,WHITE,True)
    add_text(slide,"Next review: validate completed actions against revenue recovered and target pace.",.9,5.65,10.5,.4,12,WHITE)

    bio=BytesIO(); prs.save(bio); return bio.getvalue()

# ============================================================================
# V2.8 Advanced AI Presentation Studio
# ============================================================================

PRESENTATION_FOCUS_OPTIONS = [
    "Executive Performance",
    "Target & Forecast",
    "Sales Growth",
    "Gap Diagnosis",
    "Pareto Product",
    "Pareto Migration",
    "Product Opportunity",
    "Stockout / Lost Sales Risk",
    "Inventory Health",
    "Inventory Capital at Risk",
    "Profitability",
    "Supplier Performance",
    "Category Performance",
    "Mutasi / Transfer Analysis",
    "Purchase & Replenishment",
    "Anomaly & Data Quality",
    "30-Day Action Plan",
]

PRESENTATION_DEPTHS = {
    "Executive — 7–9 slides": 9,
    "Standard — 10–14 slides": 14,
    "Deep Dive — 15–20 slides": 20,
}

SLIDE_LIBRARY = {
    "cover": {"title": "Branch Performance & Target Recovery", "objective": "Open the management story and establish the period, branch, and purpose."},
    "executive": {"title": "Executive Overview", "objective": "Summarize the most important commercial and inventory signals for management."},
    "target": {"title": "Target Performance & Forecast", "objective": "Show target achievement, pace, projection, and required sales to close the gap."},
    "sales_trend": {"title": "Sales Growth & Commercial Drivers", "objective": "Explain revenue, transaction, ATV, UPT, and month-to-month development."},
    "gap_diagnosis": {"title": "Where Are We Losing The Target?", "objective": "Diagnose the target gap using data-supported commercial and availability signals."},
    "pareto": {"title": "Pareto & Revenue Concentration", "objective": "Show Core 20, A80, Opportunity, and Long Tail concentration dynamically."},
    "pareto_migration": {"title": "Pareto Migration & Product Momentum", "objective": "Identify products moving into or out of Core/Opportunity status versus the previous month."},
    "product_opportunity": {"title": "Revenue Growth Opportunities", "objective": "Prioritize products that are most realistic to push based on revenue, growth, margin, and stock readiness."},
    "stockout_recovery": {"title": "Stockout & Revenue Recovery", "objective": "Identify proven-demand products whose availability can be recovered to protect revenue."},
    "inventory_health": {"title": "Inventory Health", "objective": "Show Normal, Slow, Dead, Overstock, No Sales, Stockout, and Negative stock exposure."},
    "inventory_capital": {"title": "Inventory Capital at Risk", "objective": "Quantify rupiah capital tied in unhealthy inventory statuses."},
    "profitability": {"title": "Profitability & Margin", "objective": "Evaluate estimated gross profit and margin development while preserving data caveats."},
    "supplier": {"title": "Supplier Capital Productivity", "objective": "Compare supplier revenue share versus inventory share and identify productive or capital-heavy suppliers."},
    "category": {"title": "Category Performance", "objective": "Compare category revenue contribution, inventory exposure, and productivity."},
    "transfer": {"title": "Mutasi / Transfer Analysis", "objective": "Explain inbound/outbound transfer flows and their relevance to branch availability and capital."},
    "purchase": {"title": "Purchase & Replenishment", "objective": "Evaluate purchase inflow and replenishment priorities relative to demand."},
    "anomaly": {"title": "Risk & Data Quality", "objective": "Surface negative stock, unresolved HPP, suspicious data, and other material anomalies."},
    "action_plan": {"title": "30-Day Action Plan", "objective": "Translate analysis into prioritized actions, owners, timing, and expected commercial impact."},
    "closing": {"title": "Management Focus", "objective": "Close with the few decisions management should carry forward from the presentation."},
}

FOCUS_TO_SLIDE = {
    "Executive Performance": "executive",
    "Target & Forecast": "target",
    "Sales Growth": "sales_trend",
    "Gap Diagnosis": "gap_diagnosis",
    "Pareto Product": "pareto",
    "Pareto Migration": "pareto_migration",
    "Product Opportunity": "product_opportunity",
    "Stockout / Lost Sales Risk": "stockout_recovery",
    "Inventory Health": "inventory_health",
    "Inventory Capital at Risk": "inventory_capital",
    "Profitability": "profitability",
    "Supplier Performance": "supplier",
    "Category Performance": "category",
    "Mutasi / Transfer Analysis": "transfer",
    "Purchase & Replenishment": "purchase",
    "Anomaly & Data Quality": "anomaly",
    "30-Day Action Plan": "action_plan",
}

AUDIENCE_DEFAULT_FOCUS = {
    "Management / Owner": [
        "Executive Performance", "Target & Forecast", "Sales Growth", "Gap Diagnosis",
        "Pareto Product", "Product Opportunity", "Stockout / Lost Sales Risk",
        "Inventory Capital at Risk", "Profitability", "Supplier Performance",
        "Category Performance", "30-Day Action Plan",
    ],
    "Buyer & Inventory": [
        "Executive Performance", "Target & Forecast", "Pareto Product", "Pareto Migration",
        "Product Opportunity", "Stockout / Lost Sales Risk", "Inventory Health",
        "Inventory Capital at Risk", "Supplier Performance", "Category Performance",
        "Mutasi / Transfer Analysis", "Purchase & Replenishment", "30-Day Action Plan",
    ],
    "Store Operations": [
        "Executive Performance", "Target & Forecast", "Sales Growth", "Gap Diagnosis",
        "Pareto Product", "Product Opportunity", "Stockout / Lost Sales Risk",
        "Inventory Health", "30-Day Action Plan",
    ],
}


def recommended_focus_for_audience(audience: str) -> list[str]:
    return list(AUDIENCE_DEFAULT_FOCUS.get(audience, AUDIENCE_DEFAULT_FOCUS["Management / Owner"]))


def _period_kpi_comparison(bundle, as_of: pd.Timestamp) -> Dict[str, Any]:
    """Current MTD versus previous-month comparable day window."""
    as_of = pd.Timestamp(as_of).normalize()
    start = as_of.to_period("M").to_timestamp()
    cur = commercial_kpis(filter_period(bundle.tx, start, as_of))
    prev_month = start - pd.DateOffset(months=1)
    prev_end_day = min(as_of.day, int(prev_month.days_in_month))
    prev_end = pd.Timestamp(prev_month.year, prev_month.month, prev_end_day)
    prev = commercial_kpis(filter_period(bundle.tx, prev_month, prev_end))
    metrics = {}
    for key in ["net_sales", "trx_count", "atv", "upt", "gross_profit", "gross_margin"]:
        c = cur.get(key)
        p = prev.get(key)
        if c is None or p is None or (isinstance(c, float) and np.isnan(c)) or (isinstance(p, float) and np.isnan(p)):
            delta = None
        elif float(p) == 0:
            delta = None
        else:
            delta = float(c) / float(p) - 1
        metrics[key] = {"current": _jsonable(c), "previous": _jsonable(p), "delta": _jsonable(delta)}
    return {
        "current_period": {"start": start.strftime("%Y-%m-%d"), "end": as_of.strftime("%Y-%m-%d")},
        "previous_comparable": {"start": prev_month.strftime("%Y-%m-%d"), "end": prev_end.strftime("%Y-%m-%d")},
        "metrics": metrics,
    }


def _pareto_migration_context(bundle, as_of: pd.Timestamp, current_inventory: pd.DataFrame, current_pareto: pd.DataFrame) -> Dict[str, Any]:
    as_of = pd.Timestamp(as_of).normalize()
    cur_start = as_of.to_period("M").to_timestamp()
    prev_end = cur_start - pd.Timedelta(days=1)
    if prev_end.year != as_of.year and prev_end < bundle.min_date:
        return {"summary": {}, "top_changes": []}
    prev_start = prev_end.to_period("M").to_timestamp()
    try:
        prev_inv = inventory_health(bundle.opening, bundle.tx, prev_end, bundle.purchases)
        prev_p = pareto_products(bundle.tx, prev_inv, prev_start, prev_end)
    except Exception:
        return {"summary": {}, "top_changes": []}
    if prev_p.empty or current_pareto.empty:
        return {"summary": {}, "top_changes": []}
    a = prev_p[["sku", "pareto_group", "revenue"]].rename(columns={"pareto_group": "previous_group", "revenue": "previous_revenue"})
    b = current_pareto[["sku", "pareto_group", "revenue", "nama_barang", "supplier"]].rename(columns={"pareto_group": "current_group", "revenue": "current_revenue"})
    m = b.merge(a, on="sku", how="outer")
    m["previous_group"] = m["previous_group"].fillna("NEW/NO_SALES")
    m["current_group"] = m["current_group"].fillna("NO_CURRENT_SALES")
    m["previous_revenue"] = m["previous_revenue"].fillna(0.0)
    m["current_revenue"] = m["current_revenue"].fillna(0.0)
    m["revenue_change"] = m["current_revenue"] - m["previous_revenue"]
    changed = m[m["previous_group"].ne(m["current_group"])].copy()
    summary = changed.groupby(["previous_group", "current_group"]).size().rename("sku_count").reset_index().sort_values("sku_count", ascending=False)
    top = changed.sort_values("revenue_change", ascending=False)
    return {
        "summary": _records(summary, ["previous_group", "current_group", "sku_count"], 20),
        "top_changes": _records(top, ["sku", "nama_barang", "supplier", "previous_group", "current_group", "previous_revenue", "current_revenue", "revenue_change"], 15),
    }


def advanced_presentation_context(bundle, as_of: pd.Timestamp, location: str = "IDK-ATP") -> Dict[str, Any]:
    """Richer aggregated fact-pack for dynamic presentations; still excludes raw transaction rows."""
    ctx = presentation_context(bundle, as_of, location)
    as_of = pd.Timestamp(as_of).normalize()
    start = as_of.to_period("M").to_timestamp()
    tx_period = filter_period(bundle.tx, start, as_of)
    inv = inventory_health(bundle.opening, bundle.tx, as_of, bundle.purchases)
    p = opportunity_scoring(pareto_products(bundle.tx, inv, start, as_of), bundle.tx, as_of)

    # Daily sales trend.
    daily = tx_period[tx_period["movement"].isin(["SALE", "SALES_RETURN"])].groupby("date", as_index=False).agg(
        net_sales=("net_sales_value", "sum"),
        net_qty=("net_sales_qty", "sum"),
    )
    trx_daily = tx_period[tx_period["movement"].eq("SALE")].groupby("date")["kd_trx"].nunique().rename("trx_count")
    if not daily.empty:
        daily = daily.merge(trx_daily, on="date", how="left")
        daily["trx_count"] = daily["trx_count"].fillna(0).astype(int)
        daily["atv"] = daily["net_sales"] / daily["trx_count"].replace(0, np.nan)

    # Inventory capital by status.
    capital = inv.copy()
    capital["positive_stock_value"] = capital["current_stock_value"].clip(lower=0)
    inv_capital = capital.groupby("inventory_status", as_index=False).agg(
        sku_count=("sku", "nunique"),
        stock_qty=("current_stock", "sum"),
        stock_value=("positive_stock_value", "sum"),
    ).sort_values("stock_value", ascending=False)

    # Stockout / availability recovery candidates: proven revenue in the current period but no stock at as-of.
    stockout = p[p["current_stock"].le(0)].copy() if not p.empty and "current_stock" in p.columns else pd.DataFrame()
    if not stockout.empty:
        elapsed = max(1, as_of.day)
        stockout["recent_daily_revenue"] = stockout["revenue"] / elapsed
        stockout = stockout.sort_values(["revenue", "opportunity_score"], ascending=[False, False])

    # Purchase and transfer flows by supplier.
    movement = tx_period.copy()
    movement_summary = movement.groupby("movement", as_index=False).agg(
        stock_in=("stock_in", "sum"),
        stock_out=("stock_out", "sum"),
        movement_value=("subtotal", "sum"),
        transaction_count=("kd_trx", "nunique"),
    ).sort_values("movement_value", ascending=False)

    purchase_rows = movement[movement["movement"].isin(["PURCHASE", "PRE_RECEIVE"])].copy()
    purchase_by_supplier = purchase_rows.groupby("supplier", as_index=False).agg(
        purchase_qty=("stock_in", "sum"),
        purchase_value=("subtotal", "sum"),
        purchase_trx=("kd_trx", "nunique"),
    ).sort_values("purchase_value", ascending=False) if not purchase_rows.empty else pd.DataFrame()

    transfer_rows = movement[movement["movement"].isin(["TRANSFER_IN", "TRANSFER_OUT", "TRANSIT_IN", "TRANSIT_OUT"])].copy()
    transfer_by_supplier = transfer_rows.groupby(["movement", "supplier"], as_index=False).agg(
        qty_in=("stock_in", "sum"),
        qty_out=("stock_out", "sum"),
        transfer_value=("subtotal", "sum"),
    ).sort_values("transfer_value", ascending=False) if not transfer_rows.empty else pd.DataFrame()

    # Profitability leaders.
    profit_products = p.copy()
    if not profit_products.empty and "gross_profit" in profit_products.columns:
        profit_products = profit_products[profit_products["gross_profit"].notna()].sort_values("gross_profit", ascending=False)

    # Supplier/category with a few more rows for slide planning.
    supplier = revenue_inventory_matrix(bundle.tx, inv, start, as_of, "supplier")
    category = revenue_inventory_matrix(bundle.tx, inv, start, as_of, "subdept")

    ctx.update({
        "commercial_comparison": _period_kpi_comparison(bundle, as_of),
        "daily_sales": _records(daily, ["date", "net_sales", "net_qty", "trx_count", "atv"], 40),
        "inventory_capital": _records(inv_capital, ["inventory_status", "sku_count", "stock_qty", "stock_value"], 20),
        "stockout_recovery": _records(stockout, ["sku", "nama_barang", "supplier", "pareto_group", "revenue", "recent_daily_revenue", "growth_30d", "gross_margin", "current_stock", "inventory_status", "opportunity_score"], 20),
        "movement_summary": _records(movement_summary, ["movement", "stock_in", "stock_out", "movement_value", "transaction_count"], 30),
        "purchase_by_supplier": _records(purchase_by_supplier, ["supplier", "purchase_qty", "purchase_value", "purchase_trx"], 15),
        "transfer_by_supplier": _records(transfer_by_supplier, ["movement", "supplier", "qty_in", "qty_out", "transfer_value"], 20),
        "profit_products": _records(profit_products, ["sku", "nama_barang", "supplier", "revenue", "gross_profit", "gross_margin", "current_stock", "inventory_status"], 15),
        "top_suppliers": _records(supplier, ["supplier", "revenue", "inventory_value", "revenue_share", "inventory_share", "productivity_index"], 15),
        "top_categories": _records(category, ["subdept", "revenue", "inventory_value", "revenue_share", "inventory_share", "productivity_index"], 15),
        "pareto_migration": _pareto_migration_context(bundle, as_of, inv, p),
    })
    return ctx


def _provider_text(api_key: str, model: str, prompt: str, provider: str) -> str:
    provider_key = (provider or "openai").strip().lower()
    if provider_key == "gemini":
        if _module_available("google.genai"):
            return _call_gemini_sdk(api_key, model, prompt)
        return _call_gemini_http(api_key, model, prompt)
    if provider_key != "openai":
        raise ValueError(f"AI provider tidak dikenal: {provider}")
    try:
        from openai import OpenAI
    except ImportError:
        return _call_openai_responses_http(api_key, model, prompt)
    try:
        client = OpenAI(api_key=api_key, timeout=120.0)
        response = client.responses.create(model=model, input=prompt)
        return response.output_text
    except Exception as exc:
        raise RuntimeError(f"OpenAI API request gagal: {exc}") from exc


def _slide_limit(depth: str) -> int:
    return int(PRESENTATION_DEPTHS.get(depth, 14))


def build_recommended_slide_plan(focus_areas: list[str], depth: str, audience: str) -> list[dict]:
    """Deterministic plan used as fallback and as the baseline the AI may improve."""
    focus = [x for x in (focus_areas or []) if x in FOCUS_TO_SLIDE]
    if not focus:
        focus = recommended_focus_for_audience(audience)
    middle = []
    for name in focus:
        stype = FOCUS_TO_SLIDE[name]
        if stype not in middle:
            middle.append(stype)
    # Always establish the story and end with action/closing.
    ordered = ["cover"]
    if "executive" not in middle:
        ordered.append("executive")
    ordered.extend(middle)
    if "action_plan" not in ordered:
        ordered.append("action_plan")
    ordered.append("closing")
    # Remove duplicates while preserving order.
    uniq = []
    for x in ordered:
        if x not in uniq:
            uniq.append(x)
    limit = _slide_limit(depth)
    if len(uniq) > limit:
        # Preserve cover, executive, action, closing; trim lower-priority middle slides.
        mandatory = [x for x in ["cover", "executive"] if x in uniq]
        ending = [x for x in ["action_plan", "closing"] if x in uniq]
        room = max(0, limit - len(mandatory) - len(ending))
        middle_trim = [x for x in uniq if x not in mandatory + ending][:room]
        uniq = mandatory + middle_trim + ending
    plan = []
    for idx, stype in enumerate(uniq, start=1):
        meta = SLIDE_LIBRARY[stype]
        plan.append({
            "include": True,
            "order": idx,
            "slide_type": stype,
            "title": meta["title"],
            "objective": meta["objective"],
            "emphasis": "",
        })
    return plan


def normalize_slide_plan(plan: list[dict], depth: str | None = None) -> list[dict]:
    valid = []
    for i, row in enumerate(plan or [], start=1):
        if not isinstance(row, dict):
            continue
        stype = str(row.get("slide_type", "")).strip()
        if stype not in SLIDE_LIBRARY:
            continue
        include = bool(row.get("include", True))
        if not include:
            continue
        try:
            order = int(row.get("order", i))
        except Exception:
            order = i
        valid.append({
            "include": True,
            "order": order,
            "slide_type": stype,
            "title": str(row.get("title") or SLIDE_LIBRARY[stype]["title"]).strip(),
            "objective": str(row.get("objective") or SLIDE_LIBRARY[stype]["objective"]).strip(),
            "emphasis": str(row.get("emphasis") or "").strip(),
        })
    valid.sort(key=lambda x: x["order"])
    if depth:
        valid = valid[:_slide_limit(depth)]
    for idx, row in enumerate(valid, start=1):
        row["order"] = idx
    return valid


def generate_ai_slide_plan(api_key: str, model: str, context: Dict[str, Any], *, audience: str, language: str,
                           objective: str, focus_areas: list[str], additional_points: str, depth: str,
                           provider: str = "openai") -> list[dict]:
    """Ask the provider to architect the slide story before presentation generation."""
    if not api_key:
        return build_recommended_slide_plan(focus_areas, depth, audience)
    baseline = build_recommended_slide_plan(focus_areas, depth, audience)
    allowed = {k: v for k, v in SLIDE_LIBRARY.items()}
    lang_instruction = "Bahasa Indonesia yang profesional" if language == "Bahasa Indonesia" else "professional English"
    prompt = f"""
You are the Lead Retail Business Analyst and presentation architect for INDOKIDS.
Design a management presentation plan BEFORE writing slide content.
Audience: {audience}
Language: {lang_instruction}
Depth: {depth} (maximum {_slide_limit(depth)} slides)

BUSINESS OBJECTIVE:
{objective or 'Evaluate branch performance and determine the most practical actions to achieve the monthly target.'}

FOCUS AREAS SELECTED BY USER:
{json.dumps(focus_areas, ensure_ascii=False)}

MANDATORY USER POINTS:
{additional_points or 'None.'}

RULES:
- Use only slide types from ALLOWED_SLIDE_LIBRARY.
- Cover, Executive Overview, 30-Day Action Plan, and Management Focus should normally be preserved unless clearly unnecessary.
- Build a coherent story: Monitor -> Diagnose -> Opportunity/Risk -> Act.
- Do not add a topic that cannot be supported by CONTEXT.
- Avoid redundant slides.
- Prefer decision-useful slides over decorative slides.
- Keep within the requested slide count.
- Return ONLY JSON.

Return shape:
{{"slides":[{{"include":true,"order":1,"slide_type":"cover","title":"...","objective":"...","emphasis":"..."}}]}}

ALLOWED_SLIDE_LIBRARY:
{json.dumps(allowed, ensure_ascii=False)}

BASELINE_PLAN:
{json.dumps(baseline, ensure_ascii=False)}

CONTEXT SUMMARY:
{json.dumps(context, ensure_ascii=False, default=_jsonable)}
""".strip()
    raw = _provider_text(api_key, model, prompt, provider)
    data = _extract_json(raw)
    plan = normalize_slide_plan(data.get("slides", []), depth)
    return plan or baseline


def generate_dynamic_presentation_content(api_key: str, model: str, context: Dict[str, Any], slide_plan: list[dict], *,
                                          audience: str, language: str, objective: str, additional_points: str,
                                          provider: str = "openai") -> Dict[str, Any]:
    """Generate slide-specific insight while keeping all numbers grounded in the application fact-pack."""
    if not api_key:
        raise ValueError("API key belum diisi.")
    plan = normalize_slide_plan(slide_plan)
    lang_instruction = "Bahasa Indonesia yang profesional, tajam, dan mudah dipresentasikan" if language == "Bahasa Indonesia" else "professional, sharp, presentation-ready English"
    output_shape = {
        "presentation_title": "string",
        "executive_takeaway": "string",
        "slides": [{
            "slide_type": "one of requested slide types",
            "title": "string",
            "headline": "one decisive sentence",
            "bullets": ["2-5 concise evidence-based bullets"],
            "risk_note": "optional string",
            "action_note": "optional string",
        }],
        "recommended_actions": [{
            "priority": 1,
            "action": "string",
            "why": "string",
            "expected_impact": "string",
            "owner": "string",
            "timing": "string",
        }],
        "closing_message": "string",
    }
    prompt = f"""
You are the Lead Retail Business Analyst for INDOKIDS, preparing a management presentation.
Audience: {audience}
Write in: {lang_instruction}.

YOUR JOB IS NOT TO DESCRIBE CHARTS. Determine:
1. What happened?
2. Why is it commercially important?
3. Where is the largest opportunity or risk?
4. What should management do next?
5. Which SKU, category, supplier, or inventory action should be prioritized?

STRICT EVIDENCE RULES:
- Use ONLY facts and numbers contained in CONTEXT.
- Never invent revenue, margin, target, causal explanations, SKU facts, or financial impact.
- Separate FACT from DATA-SUPPORTED INTERPRETATION and HYPOTHESIS.
- If causality is not proven, explicitly frame it as a hypothesis or recommended investigation.
- HPP/Gross Profit are estimates when present.
- Pareto is dynamic. Core 20% is NOT automatically 80% of revenue.
- The commercial objective is to achieve the monthly target without sacrificing margin or creating unhealthy inventory.
- Every material problem should identify the metric/entity and a practical next action.
- Do not mention that you are AI.
- Keep slide bullets concise enough to fit a professional PowerPoint.

PRESENTATION OBJECTIVE:
{objective}

MANDATORY USER BRIEF:
{additional_points or 'No additional mandatory points.'}

APPROVED SLIDE PLAN (follow this order and these slide types):
{json.dumps(plan, ensure_ascii=False)}

Return ONLY valid JSON matching:
{json.dumps(output_shape, ensure_ascii=False)}

CONTEXT:
{json.dumps(context, ensure_ascii=False, default=_jsonable)}
""".strip()
    raw = _provider_text(api_key, model, prompt, provider)
    data = _extract_json(raw)

    slide_map = {}
    for item in data.get("slides", []) if isinstance(data.get("slides", []), list) else []:
        if isinstance(item, dict) and item.get("slide_type") in SLIDE_LIBRARY:
            slide_map[item["slide_type"]] = {
                "slide_type": item["slide_type"],
                "title": str(item.get("title") or SLIDE_LIBRARY[item["slide_type"]]["title"]),
                "headline": str(item.get("headline") or ""),
                "bullets": [str(x) for x in item.get("bullets", []) if str(x).strip()][:6],
                "risk_note": str(item.get("risk_note") or ""),
                "action_note": str(item.get("action_note") or ""),
            }
    slides = []
    for row in plan:
        stype = row["slide_type"]
        item = slide_map.get(stype, {
            "slide_type": stype,
            "title": row["title"],
            "headline": "",
            "bullets": [],
            "risk_note": "",
            "action_note": "",
        })
        item["title"] = row.get("title") or item["title"]
        slides.append(item)
    data["slides"] = slides

    actions = data.get("recommended_actions", [])
    norm = []
    if isinstance(actions, list):
        for i, row in enumerate(actions[:10], 1):
            if not isinstance(row, dict):
                continue
            norm.append({
                "priority": row.get("priority", i),
                "action": str(row.get("action", "")),
                "why": str(row.get("why", "")),
                "expected_impact": str(row.get("expected_impact", "")),
                "owner": str(row.get("owner", "Management")),
                "timing": str(row.get("timing", "")),
            })
    data["recommended_actions"] = norm
    data["presentation_title"] = str(data.get("presentation_title") or "Branch Performance & Target Recovery")
    data["closing_message"] = str(data.get("closing_message") or "Protect winners. Grow potentials. Fix availability. Reduce waste.")
    return data


def _slide_content(ai: Dict[str, Any], slide_type: str) -> Dict[str, Any]:
    for row in ai.get("slides", []) or []:
        if isinstance(row, dict) and row.get("slide_type") == slide_type:
            return row
    return {"title": SLIDE_LIBRARY.get(slide_type, {}).get("title", slide_type), "headline": "", "bullets": [], "risk_note": "", "action_note": ""}


def build_dynamic_pptx(context: Dict[str, Any], ai: Dict[str, Any], slide_plan: list[dict], *, branch: str,
                       as_of: pd.Timestamp, audience: str, objective: str = "") -> bytes:
    """Build a variable slide deck from an approved plan. All charts/tables use engine-calculated context."""
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("Package 'python-pptx' belum terinstall.") from exc

    plan = normalize_slide_plan(slide_plan)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5); blank = prs.slide_layouts[6]
    BLUE=RGBColor(47,104,176); ORANGE=RGBColor(246,139,36); DARK=RGBColor(31,41,55); MUTED=RGBColor(100,116,139)
    LIGHT=RGBColor(244,247,251); WHITE=RGBColor(255,255,255); GREEN=RGBColor(22,163,74); RED=RGBColor(220,38,38)

    def add_rect(slide,x,y,w,h,fill,radius=False):
        from pptx.enum.shapes import MSO_SHAPE
        shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h)); shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.fill.background(); return shp
    def add_text(slide,text,x,y,w,h,size=18,color=DARK,bold=False,align=None):
        box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True; p=tf.paragraphs[0]; p.text=str(text or ""); p.font.name="Aptos"; p.font.size=Pt(size); p.font.bold=bold; p.font.color.rgb=color; p.alignment=align if align is not None else p.alignment; return box
    def header(slide,title,subtitle,num):
        add_text(slide,title,.7,.32,11.55,.52,24,DARK,True); add_text(slide,subtitle or "",.72,.88,11.6,.32,10.2,MUTED); add_rect(slide,.7,1.23,1,.055,ORANGE); add_text(slide,str(num),12.15,.38,.45,.28,9,MUTED,False,PP_ALIGN.RIGHT)
    def footer(slide): add_text(slide,f"INDOKIDS · {branch} · Data s.d. {pd.Timestamp(as_of):%d %b %Y}",.7,7.12,11.8,.22,8.5,MUTED)
    def bullets(slide,items,x,y,w,h,size=14):
        box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
        for i,item in enumerate(items or []):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text="• "+str(item); p.font.name="Aptos"; p.font.size=Pt(size); p.font.color.rgb=DARK; p.space_after=Pt(8)
        return box
    def metric(slide,x,y,w,title,value,note=""):
        add_rect(slide,x,y,w,1.02,LIGHT,True); add_text(slide,title,x+.15,y+.11,w-.3,.2,9,MUTED,True); add_text(slide,value,x+.15,y+.39,w-.3,.3,18,DARK,True); add_text(slide,note,x+.15,y+.76,w-.3,.18,8.2,MUTED) if note else None
    def insight_side(slide, content, x=8.0, y=1.6, w=4.5, h=4.9):
        headline=content.get("headline",""); b=content.get("bullets",[]); risk=content.get("risk_note",""); action=content.get("action_note","")
        if headline: add_text(slide,headline,x,y,w,.7,15,BLUE,True)
        bullets(slide,b,x,y+.82,w,h-1.45,13)
        if risk: add_text(slide,"Risk: "+risk,x,y+h-.55,w,.25,9.5,RED,True)
        elif action: add_text(slide,"Action: "+action,x,y+h-.55,w,.25,9.5,GREEN,True)
    def table_slide(slide, rows, headers, keys, x=.65,y=1.55,w=12.0,h=4.95, max_rows=8):
        rows=(rows or [])[:max_rows]
        if not rows: add_text(slide,"Tidak ada data yang cukup untuk tabel ini.",.8,2.2,11.5,.5,16,MUTED); return
        table=slide.shapes.add_table(len(rows)+1,len(headers),Inches(x),Inches(y),Inches(w),Inches(h)).table
        widths=[w/len(headers)]*len(headers)
        for i,ww in enumerate(widths): table.columns[i].width=Inches(ww)
        for c,hdr in enumerate(headers):
            cell=table.cell(0,c); cell.text=hdr; cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
            for p in cell.text_frame.paragraphs: p.font.color.rgb=WHITE; p.font.bold=True; p.font.size=Pt(9)
        for ri,row in enumerate(rows,1):
            for ci,key in enumerate(keys):
                val=row.get(key,"")
                if key in {"revenue","inventory_value","stock_value","purchase_value","transfer_value","gross_profit","current_revenue","previous_revenue","revenue_change"}: val=rupiah(val)
                elif key in {"revenue_share","inventory_share","gross_margin","growth_30d"}: val=pct(val)
                elif isinstance(val,float): val=f"{val:,.1f}".replace(",",".")
                cell=table.cell(ri,ci); cell.text=str(val); cell.fill.solid(); cell.fill.fore_color.rgb=WHITE if ri%2 else LIGHT
                for p in cell.text_frame.paragraphs: p.font.size=Pt(8.1); p.font.color.rgb=DARK

    k=context.get("kpi",{}); t=context.get("target") or {}; par=context.get("pareto",{}); monthly=context.get("monthly",[])
    total_slides=len(plan)
    for idx,row in enumerate(plan,1):
        stype=row["slide_type"]; content=_slide_content(ai,stype); title=row.get("title") or content.get("title") or SLIDE_LIBRARY[stype]["title"]
        if stype=="cover":
            s=prs.slides.add_slide(blank); add_rect(s,0,0,13.333,7.5,BLUE); add_rect(s,.72,1.35,.12,2.45,ORANGE)
            add_text(s,ai.get("presentation_title") or title,1.1,1.35,10.9,1.35,31,WHITE,True); add_text(s,f"{branch} · {audience}",1.1,2.95,9.5,.45,18,WHITE)
            add_text(s,f"Data sampai {pd.Timestamp(as_of):%d %B %Y}",1.1,3.5,8,.32,11,WHITE); add_text(s,objective or "Monitor · Diagnose · Act",1.1,5.75,10.8,.6,13,WHITE,True); continue
        if stype=="closing":
            s=prs.slides.add_slide(blank); add_rect(s,0,0,13.333,7.5,DARK); add_text(s,title,.85,1.0,6,.6,29,WHITE,True); add_rect(s,.87,1.83,1,.06,ORANGE); add_text(s,ai.get("closing_message") or content.get("headline"),.87,2.35,10.9,1.45,24,WHITE,True); bullets(s,content.get("bullets",[]),.9,4.25,10.8,1.6,13); continue
        s=prs.slides.add_slide(blank); header(s,title,row.get("objective",""),idx)
        if stype=="executive":
            metric(s,.7,1.52,2.82,"NET SALES MTD",rupiah(k.get("net_sales"))); metric(s,3.7,1.52,2.82,"TARGET",rupiah(t.get("target")) if t else "-"); metric(s,6.7,1.52,2.82,"ACHIEVEMENT",pct(t.get("achievement")) if t else "-"); metric(s,9.7,1.52,2.9,"PROJECTED GAP",rupiah(t.get("projected_gap")) if t else "-")
            metric(s,.7,2.72,2.82,"TRANSACTIONS",f"{int(k.get('trx_count',0)):,}".replace(",",".")); metric(s,3.7,2.72,2.82,"ATV",rupiah(k.get("atv"))); metric(s,6.7,2.72,2.82,"EST. GROSS PROFIT",rupiah(k.get("gross_profit"))); metric(s,9.7,2.72,2.9,"GROSS MARGIN",pct(k.get("gross_margin")))
            add_text(s,content.get("headline","Management readout"),.75,4.1,11.8,.45,15,BLUE,True); bullets(s,content.get("bullets",[]),.8,4.65,11.7,1.85,14)
        elif stype in {"target","sales_trend"}:
            if monthly:
                cd=CategoryChartData(); cd.categories=[pd.Timestamp(r["month"]).strftime("%b") for r in monthly]; cd.add_series("Net Sales",[float(r.get("net_sales") or 0) for r in monthly])
                if stype=="target": cd.add_series("Target",[float(r.get("target_omzet") or 0) for r in monthly])
                chart=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED if stype=="target" else XL_CHART_TYPE.LINE_MARKERS,Inches(.75),Inches(1.55),Inches(7.0),Inches(4.9),cd).chart; chart.has_legend=True; chart.legend.position=XL_LEGEND_POSITION.BOTTOM; chart.value_axis.tick_labels.number_format='0,,"M"'
            insight_side(s,content,8.05,1.55,4.5,4.95)
        elif stype=="gap_diagnosis":
            metric(s,.75,1.55,2.75,"ACTUAL",rupiah(t.get("actual")) if t else "-"); metric(s,3.7,1.55,2.75,"TARGET GAP",rupiah(t.get("gap")) if t else "-"); metric(s,6.65,1.55,2.75,"PROJECTED",rupiah(t.get("projected_month_end")) if t else "-"); metric(s,9.6,1.55,2.75,"REQUIRED / DAY",rupiah(t.get("required_daily_sales")) if t else "-")
            comp=context.get("commercial_comparison",{}).get("metrics",{}); labels=[("Sales",comp.get("net_sales",{})),("Transactions",comp.get("trx_count",{})),("ATV",comp.get("atv",{})),("UPT",comp.get("upt",{}))]
            y=3.05
            for lab,m in labels:
                delta=m.get("delta"); txt=pct(delta) if delta is not None else "n/a"; add_text(s,lab,.9,y,2.2,.3,11,DARK,True); add_text(s,txt,3.0,y,1.3,.3,12,GREEN if (delta or 0)>=0 else RED,True); y+=.58
            insight_side(s,content,5.0,3.0,7.3,3.1)
        elif stype=="pareto":
            cd=CategoryChartData(); cd.categories=["Core 20","Opportunity","Long Tail"]; cd.add_series("Revenue Share",[par.get("core_share",0)*100,par.get("opportunity_share",0)*100,max(0,1-par.get("core_share",0)-par.get("opportunity_share",0))*100]); ch=s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,Inches(.8),Inches(1.65),Inches(5.25),Inches(3.8),cd).chart; ch.has_legend=False; ch.value_axis.maximum_scale=100; ch.value_axis.tick_labels.number_format='0"%"'
            metric(s,.85,5.65,2.35,"CORE SHARE",pct(par.get("core_share"))); metric(s,3.4,5.65,2.35,"A80 SKU",f"{par.get('a80_sku',0):,}".replace(",",".")); insight_side(s,content,6.4,1.65,6.0,4.95)
        elif stype=="pareto_migration":
            rows=context.get("pareto_migration",{}).get("top_changes",[]); table_slide(s,rows,["SKU","Product","Previous","Current","Revenue Δ"],["sku","nama_barang","previous_group","current_group","revenue_change"],.65,1.55,7.3,4.95,8); insight_side(s,content,8.25,1.55,4.1,4.95)
        elif stype=="product_opportunity":
            rows=context.get("top_focus_products",[]); table_slide(s,rows,["SKU","Product","Pareto","Revenue","Growth","Stock","Status","Action"],["sku","nama_barang","pareto_group","revenue","growth_30d","current_stock","inventory_status","recommended_action"],.55,1.5,12.25,4.65,8); add_text(s,content.get("headline",""),.7,6.35,12,.32,12,BLUE,True)
        elif stype=="stockout_recovery":
            rows=context.get("stockout_recovery",[]); table_slide(s,rows,["SKU","Product","Pareto","Revenue","Daily Rev","Growth","Status"],["sku","nama_barang","pareto_group","revenue","recent_daily_revenue","growth_30d","inventory_status"],.6,1.55,8.0,4.85,8); insight_side(s,content,8.85,1.55,3.6,4.85)
        elif stype in {"inventory_health","inventory_capital"}:
            rows=context.get("inventory_capital" if stype=="inventory_capital" else "inventory_summary",[])
            if rows:
                cd=CategoryChartData(); cd.categories=[str(r.get("inventory_status","")) for r in rows]; key="stock_value" if stype=="inventory_capital" else "sku_count"; cd.add_series("Stock Value" if key=="stock_value" else "SKU",[float(r.get(key) or 0) for r in rows]); ch=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(.75),Inches(1.6),Inches(7.0),Inches(4.7),cd).chart; ch.has_legend=False; ch.value_axis.tick_labels.number_format='0,,"M"' if key=="stock_value" else '0'
            insight_side(s,content,8.05,1.6,4.45,4.8)
        elif stype=="profitability":
            metric(s,.75,1.5,2.8,"EST. GROSS PROFIT",rupiah(k.get("gross_profit"))); metric(s,3.75,1.5,2.8,"GROSS MARGIN",pct(k.get("gross_margin"))); rows=context.get("profit_products",[]); table_slide(s,rows,["SKU","Product","Revenue","Gross Profit","Margin"],["sku","nama_barang","revenue","gross_profit","gross_margin"],.65,2.85,7.4,3.35,6); insight_side(s,content,8.35,2.0,4.0,4.5)
        elif stype in {"supplier","category"}:
            rows=context.get("top_suppliers" if stype=="supplier" else "top_categories",[]); dim="supplier" if stype=="supplier" else "subdept"; hdr="Supplier" if stype=="supplier" else "Category"; table_slide(s,rows,[hdr,"Revenue","Inventory","Rev Share","Inv Share","Productivity"],[dim,"revenue","inventory_value","revenue_share","inventory_share","productivity_index"],.55,1.55,8.15,4.9,9); insight_side(s,content,8.95,1.55,3.45,4.9)
        elif stype=="transfer":
            rows=context.get("transfer_by_supplier",[]); table_slide(s,rows,["Movement","Supplier","Qty In","Qty Out","Value"],["movement","supplier","qty_in","qty_out","transfer_value"],.6,1.55,8.1,4.85,9); insight_side(s,content,8.95,1.55,3.4,4.85)
        elif stype=="purchase":
            rows=context.get("purchase_by_supplier",[]); table_slide(s,rows,["Supplier","Qty","Purchase Value","Trx"],["supplier","purchase_qty","purchase_value","purchase_trx"],.6,1.55,8.1,4.85,9); insight_side(s,content,8.95,1.55,3.4,4.85)
        elif stype=="anomaly":
            rows=[{"issue":k,"count":v} for k,v in (context.get("anomalies") or {}).items()]; table_slide(s,rows,["Issue","Rows"],["issue","count"],.75,1.6,5.2,4.6,12); insight_side(s,content,6.35,1.6,5.95,4.7)
        elif stype=="action_plan":
            actions=ai.get("recommended_actions",[])[:8]; headers=["#","Action","Why","Impact","Owner","Timing"]; keys=["priority","action","why","expected_impact","owner","timing"]; table_slide(s,actions,headers,keys,.45,1.48,12.45,5.2,8)
        else:
            insight_side(s,content,.8,1.6,11.8,4.9)
        footer(s)
    bio=BytesIO(); prs.save(bio); return bio.getvalue()
